from pathlib import Path
from PIL import Image, ImageEnhance, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
SRC = ROOT / "tmp" / "screens"
PROC = ROOT / "tmp" / "processed"
OUT = ROOT / "output"
OUT.mkdir(parents=True, exist_ok=True)
PROC.mkdir(parents=True, exist_ok=True)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(10, 17, 38)
NAVY_2 = RGBColor(18, 27, 50)
SLATE = RGBColor(31, 41, 55)
TEXT = RGBColor(248, 250, 252)
MUTED = RGBColor(203, 213, 225)
CYAN = RGBColor(48, 211, 255)
TEAL = RGBColor(34, 197, 160)
LIME = RGBColor(163, 230, 53)
AMBER = RGBColor(251, 191, 36)
ROSE = RGBColor(244, 63, 94)
CARD = RGBColor(16, 24, 43)
CARD_2 = RGBColor(22, 32, 56)
WHITE = RGBColor(255, 255, 255)


def hex_rgb(value):
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def prep_image(src_name, dest_name, size, crop=None, brightness=1.0, contrast=1.0):
    src = SRC / src_name
    dst = PROC / dest_name
    with Image.open(src) as img:
        if crop:
            img = img.crop(crop)
        img = ImageOps.fit(img, size, method=Image.Resampling.LANCZOS)
        if brightness != 1.0:
            img = ImageEnhance.Brightness(img).enhance(brightness)
        if contrast != 1.0:
            img = ImageEnhance.Contrast(img).enhance(contrast)
        img.save(dst)
    return dst


def add_bg(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_full_rect(slide, left, top, width, height, fill, line=None, transparency=0):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if transparency:
        shape.fill.transparency = transparency
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_round_card(slide, left, top, width, height, fill=CARD, line=SLATE, radius=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.1)
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font="Aptos", italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return box


def add_multiline_text(slide, left, top, width, height, lines, font_size=18, color=TEXT, bold=False, spacing=1.05, font="Aptos"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font
        p.line_spacing = spacing
    return box


def add_image(slide, image_path, left, top, width, height, rounded=False):
    if rounded:
        frame = add_round_card(slide, left, top, width, height, fill=CARD_2, line=SLATE)
        pic = slide.shapes.add_picture(str(image_path), left + Pt(5), top + Pt(5), width - Pt(10), height - Pt(10))
        return frame, pic
    return slide.shapes.add_picture(str(image_path), left, top, width, height)


def add_footer(slide, num, total):
    add_textbox(slide, Inches(11.8), Inches(7.05), Inches(1.0), Inches(0.2), f"{num}/{total}", font_size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_pill(slide, left, top, width, text, fill, font_size=13):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.36))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.text = text
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.size = Pt(font_size)
    r.font.bold = True
    r.font.color.rgb = NAVY
    r.font.name = "Aptos"
    return shape


def add_big_title(slide, title, subtitle, eyebrow=None):
    if eyebrow:
        add_pill(slide, Inches(0.7), Inches(0.55), Inches(2.1), eyebrow, TEAL, 12)
    add_textbox(slide, Inches(0.7), Inches(1.0), Inches(6.3), Inches(1.1), title, font_size=27, bold=True, color=TEXT)
    add_textbox(slide, Inches(0.7), Inches(2.0), Inches(6.0), Inches(0.75), subtitle, font_size=15, color=MUTED)


def add_bullets(slide, left, top, width, height, bullets, title=None, accent=CYAN):
    add_round_card(slide, left, top, width, height, fill=CARD, line=SLATE)
    if title:
        add_textbox(slide, left + Inches(0.2), top + Inches(0.12), width - Inches(0.4), Inches(0.28), title, font_size=16, bold=True, color=accent)
    box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.42), width - Inches(0.4), height - Inches(0.5))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT
        p.font.name = "Aptos"
        p.space_after = Pt(10)
        p.bullet = True
    return box


def add_card_with_label(slide, left, top, width, height, label, body, accent=CYAN):
    add_round_card(slide, left, top, width, height, fill=CARD_2, line=SLATE)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.15), width - Inches(0.36), Inches(0.28), label, font_size=12, bold=True, color=accent)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.42), width - Inches(0.36), height - Inches(0.5), body, font_size=14, color=TEXT)


def add_screen_tile(slide, image_path, left, top, width, height, label, accent=CYAN):
    add_round_card(slide, left, top, width, height, fill=CARD_2, line=SLATE)
    pic = slide.shapes.add_picture(str(image_path), left + Inches(0.06), top + Inches(0.34), width - Inches(0.12), height - Inches(0.40))
    add_pill(slide, left + Inches(0.18), top + Inches(0.08), Inches(1.3), label, accent, 11)
    return pic


def add_flow_box(slide, left, top, width, height, label, detail, fill):
    shape = add_round_card(slide, left, top, width, height, fill=fill, line=fill)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(6)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = label
    r1.font.size = Pt(14)
    r1.font.bold = True
    r1.font.color.rgb = NAVY
    r1.font.name = "Aptos"
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = detail
    r2.font.size = Pt(10)
    r2.font.color.rgb = NAVY
    r2.font.name = "Aptos"
    return shape


def add_arrow(slide, left, top, width, height, fill=CYAN):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    queue_img = prep_image("queue.png", "queue_fit.png", (780, 520), brightness=1.0, contrast=1.08)
    incident_img = prep_image("incident.png", "incident_fit.png", (860, 560), crop=(0, 0, 1440, 2650), brightness=1.0, contrast=1.06)
    inputs_img = prep_image("inputs.png", "inputs_fit.png", (620, 430), brightness=1.0, contrast=1.04)
    replay_img = prep_image("replay.png", "replay_fit.png", (560, 340), brightness=1.0, contrast=1.04)
    settings_img = prep_image("settings.png", "settings_fit.png", (560, 340), brightness=1.0, contrast=1.04)
    training_img = prep_image("training.png", "training_fit.png", (560, 340), crop=(0, 0, 1440, 2300), brightness=1.0, contrast=1.04)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    add_full_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    add_full_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.18), TEAL)
    add_full_rect(slide, Inches(6.55), Inches(0.0), Inches(0.35), Inches(7.5), CYAN, transparency=0.88)
    add_full_rect(slide, Inches(11.9), Inches(-0.1), Inches(2.2), Inches(2.1), TEAL, transparency=0.88)
    add_full_rect(slide, Inches(-0.2), Inches(6.25), Inches(3.0), Inches(1.3), AMBER, transparency=0.9)

    add_pill(slide, Inches(0.7), Inches(0.45), Inches(2.65), "NEXUS v2 | STAR pitch", TEAL, 12)
    add_textbox(slide, Inches(0.7), Inches(1.0), Inches(6.0), Inches(1.1), "Situation, Task, Action, Result", font_size=28, bold=True, color=TEXT)
    add_textbox(slide, Inches(0.72), Inches(2.0), Inches(5.9), Inches(0.85), "A queue-first incident response product for business, product, engineering, and operations stakeholders.", font_size=16, color=MUTED)
    add_card_with_label(slide, Inches(0.7), Inches(3.0), Inches(5.65), Inches(1.15), "Why this matters", "Fragmented intake, unclear ownership, and hidden evidence slow incident work and make demos feel incomplete.", accent=AMBER)
    add_card_with_label(slide, Inches(0.7), Inches(4.35), Inches(5.65), Inches(1.15), "What NEXUS shows", "A visible queue, an explainable incident console, and a governed path from intake to outcome.", accent=CYAN)

    add_round_card(slide, Inches(7.1), Inches(0.78), Inches(5.45), Inches(4.55), fill=CARD_2, line=SLATE)
    slide.shapes.add_picture(str(queue_img), Inches(7.22), Inches(0.9), Inches(5.2), Inches(4.3))
    add_textbox(slide, Inches(7.28), Inches(4.98), Inches(4.8), Inches(0.3), "Live queue-first landing surface", font_size=12, color=MUTED)

    # STAR ribbon
    y = Inches(6.0)
    card_w = Inches(3.05)
    gap = Inches(0.18)
    star = [
        ("Situation", "Incident intake is fragmented across channels.", AMBER),
        ("Task", "Create one operating model with visible priority.", CYAN),
        ("Action", "Show the live product proving the workflow.", TEAL),
        ("Result", "Deliver clarity, credibility, and a path to hardening.", LIME),
    ]
    x = Inches(0.68)
    for title, body, accent in star:
        add_round_card(slide, x, y, card_w, Inches(0.92), fill=CARD, line=accent)
        add_textbox(slide, x + Inches(0.14), y + Inches(0.10), card_w - Inches(0.28), Inches(0.2), title, font_size=12, bold=True, color=accent)
        add_textbox(slide, x + Inches(0.14), y + Inches(0.31), card_w - Inches(0.28), Inches(0.42), body, font_size=10.5, color=TEXT)
        x += card_w + gap
    add_notes(slide, "Open with the problem. Explain that STAR is the speaking structure, but the product language remains enterprise-friendly. Point to the live queue screenshot as proof that this is not a mockup.")
    add_footer(slide, 1, 5)

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY_2)
    add_full_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY_2)
    add_big_title(slide, "Situation and task", "The system is designed to convert fragmented incident handling into one understandable operating model.", "01 / STAR")
    add_round_card(slide, Inches(0.7), Inches(2.35), Inches(5.8), Inches(3.85), fill=CARD, line=SLATE)
    add_textbox(slide, Inches(0.95), Inches(2.55), Inches(2.2), Inches(0.3), "Situation", font_size=16, bold=True, color=AMBER)
    add_card_with_label(slide, Inches(0.95), Inches(2.95), Inches(5.3), Inches(0.9), "Fragmented intake", "Webhook, manual reports, chat commands, streams, and batch data all arrive differently.", accent=AMBER)
    add_card_with_label(slide, Inches(0.95), Inches(4.0), Inches(5.3), Inches(0.9), "Slow triage", "Operators lose time deciding what is first, what is real, and what changed.", accent=ROSE)
    add_card_with_label(slide, Inches(0.95), Inches(5.05), Inches(5.3), Inches(0.9), "Hidden evidence", "The business cannot easily see why an incident happened or what action was taken.", accent=CYAN)

    add_round_card(slide, Inches(6.75), Inches(2.35), Inches(5.9), Inches(3.85), fill=CARD, line=SLATE)
    add_textbox(slide, Inches(7.0), Inches(2.55), Inches(2.0), Inches(0.3), "Task", font_size=16, bold=True, color=TEAL)
    add_card_with_label(slide, Inches(7.0), Inches(2.95), Inches(5.35), Inches(0.9), "Normalize everything", "Every intake path lands in the same incident envelope and queue.", accent=TEAL)
    add_card_with_label(slide, Inches(7.0), Inches(4.0), Inches(5.35), Inches(0.9), "Show the whole incident", "The console presents workflow, evidence, audit, and execution state together.", accent=CYAN)
    add_card_with_label(slide, Inches(7.0), Inches(5.05), Inches(5.35), Inches(0.9), "Make action governed", "The system is explicit about what is visible now and what still needs hardening.", accent=LIME)

    add_textbox(slide, Inches(0.82), Inches(6.45), Inches(12.0), Inches(0.4), "Supported intake channels: Webhook • Manual form • Slack-style command • Stream anomaly • Batch import", font_size=13, color=MUTED, align=PP_ALIGN.CENTER)
    add_notes(slide, "Describe the operational pain first, then the product task. This slide establishes why the product exists and why the queue-first model is the right response.")
    add_footer(slide, 2, 5)

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    add_full_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    add_big_title(slide, "Action: the live product already proves the workflow", "These screens demonstrate intake, queueing, incident analysis, and controlled execution on the same product surface.", "02 / STAR")

    add_screen_tile(slide, queue_img, Inches(0.72), Inches(2.35), Inches(3.0), Inches(2.05), "Queue", accent=AMBER)
    add_screen_tile(slide, inputs_img, Inches(0.72), Inches(4.6), Inches(3.0), Inches(1.95), "Inputs", accent=TEAL)
    add_round_card(slide, Inches(3.95), Inches(2.35), Inches(8.65), Inches(4.2), fill=CARD_2, line=SLATE)
    slide.shapes.add_picture(str(incident_img), Inches(4.07), Inches(2.48), Inches(8.4), Inches(3.9))
    add_textbox(slide, Inches(4.16), Inches(6.43), Inches(7.7), Inches(0.3), "Incident console: workflow timeline, evidence provenance, audit trail, and execution state", font_size=13, color=MUTED)

    # bottom callouts
    callouts = [
        ("Workflow timeline", CYAN),
        ("Agent handoff", TEAL),
        ("Evidence provenance", AMBER),
        ("Audit trail", LIME),
        ("Execution state", ROSE),
    ]
    cx = Inches(0.82)
    for label, accent in callouts:
        add_pill(slide, cx, Inches(6.9), Inches(2.1) if label != "Execution state" else Inches(1.8), label, accent, 11)
        cx += Inches(2.22) if label != "Execution state" else Inches(1.95)
    add_notes(slide, "Use this slide to prove the product. Walk the judge from queue to intake to the incident console, then point to evidence, audit, and the governed action path.")
    add_footer(slide, 3, 5)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY_2)
    add_full_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY_2)
    add_big_title(slide, "Result: value for every stakeholder", "The product is credible now, understandable to non-technical reviewers, and structured for the next hardening phase.", "03 / STAR")

    cards = [
        ("Business / leadership", "A visible enterprise story with control, traceability, and a clear reason to believe.", AMBER),
        ("Product", "A coherent user journey from intake to console to history, replay, training, and settings.", CYAN),
        ("Engineering", "Clear contracts and a bounded architecture that can be hardened without rewiring the UI.", TEAL),
        ("Triage / operations", "A queue-first console that shows what is next, what happened, and what action is safe.", LIME),
    ]
    card_positions = [
        (Inches(0.78), Inches(2.45)),
        (Inches(0.78), Inches(3.78)),
        (Inches(0.78), Inches(5.11)),
        (Inches(0.78), Inches(6.44)),
    ]
    for (label, body, accent), (left, top) in zip(cards, card_positions):
        add_card_with_label(slide, left, top, Inches(5.55), Inches(1.08), label, body, accent=accent)

    add_screen_tile(slide, training_img, Inches(6.6), Inches(2.45), Inches(2.95), Inches(2.0), "Training", accent=TEAL)
    add_screen_tile(slide, settings_img, Inches(9.75), Inches(2.45), Inches(2.95), Inches(2.0), "Settings", accent=AMBER)
    add_round_card(slide, Inches(6.6), Inches(4.72), Inches(6.1), Inches(1.75), fill=CARD, line=SLATE)
    add_textbox(slide, Inches(6.86), Inches(4.95), Inches(5.6), Inches(0.34), "Why this matters to the judging body", font_size=15, bold=True, color=TEXT)
    add_textbox(slide, Inches(6.86), Inches(5.32), Inches(5.6), Inches(0.85), "The product is not just a concept. It is a live, explainable workflow that already demonstrates the business story, the user experience, and the technical direction.", font_size=14, color=MUTED)
    add_round_card(slide, Inches(6.6), Inches(6.62), Inches(6.1), Inches(0.45), fill=TEAL, line=TEAL)
    add_textbox(slide, Inches(6.82), Inches(6.70), Inches(5.7), Inches(0.2), "Queue-first. Explainable. Auditable. Ready to harden.", font_size=14, bold=True, color=NAVY)
    add_notes(slide, "Shift to the stakeholder lens. Explain that this solves a business problem, creates a product story, and gives engineering a bounded path forward. The training and settings screenshots show maturity and trust.")
    add_footer(slide, 4, 5)

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    add_full_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    add_big_title(slide, "Technical design and next phase", "The current implementation is demo-ready, with clear seams for the deeper enterprise hardening work.", "04 / STAR")

    # Architecture flow
    flow_y = Inches(2.1)
    box_w = Inches(1.95)
    box_h = Inches(0.86)
    gap = Inches(0.22)
    boxes = [
        ("Intake", "Webhook, manual, chat, stream, batch", AMBER),
        ("Normalize", "One incident envelope", CYAN),
        ("Queue", "Priority, stage, ETA", TEAL),
        ("Console", "Evidence, audit, action", LIME),
        ("Replay / Training", "Learning and validation", ROSE),
        ("Settings", "Trust and readiness", AMBER),
    ]
    x = Inches(0.65)
    for idx, (label, detail, fill) in enumerate(boxes):
        add_flow_box(slide, x, flow_y, box_w, box_h, label, detail, fill)
        if idx < len(boxes) - 1:
            add_arrow(slide, x + box_w + Inches(0.03), flow_y + Inches(0.19), Inches(0.17), Inches(0.46), fill=CYAN if idx % 2 == 0 else TEAL)
        x += box_w + gap + Inches(0.17)

    add_round_card(slide, Inches(0.75), Inches(3.35), Inches(5.85), Inches(2.65), fill=CARD, line=SLATE)
    add_textbox(slide, Inches(0.98), Inches(3.55), Inches(2.8), Inches(0.3), "Current implementation", font_size=16, bold=True, color=CYAN)
    add_multiline_text(
        slide,
        Inches(1.0),
        Inches(3.95),
        Inches(5.2),
        Inches(1.7),
        [
            "FastAPI app serving the UI and versioned contracts",
            "Static frontend pages with page-specific controllers",
            "File-backed local persistence for incidents and audit",
            "Signed intake and clear queue/console seams",
        ],
        font_size=14,
        color=TEXT,
    )

    add_round_card(slide, Inches(6.9), Inches(3.35), Inches(5.8), Inches(2.65), fill=CARD, line=SLATE)
    add_textbox(slide, Inches(7.12), Inches(3.55), Inches(2.3), Inches(0.3), "Next phase", font_size=16, bold=True, color=AMBER)
    add_multiline_text(
        slide,
        Inches(7.15),
        Inches(3.95),
        Inches(5.15),
        Inches(1.7),
        [
            "Real observability integrations and evidence fusion",
            "Stronger GUARDIAN execution policy and governance",
            "Production persistence and hardened auth/tenant boundaries",
            "Further backend service decomposition and operational cleanup",
        ],
        font_size=14,
        color=TEXT,
    )

    add_round_card(slide, Inches(0.75), Inches(6.18), Inches(11.95), Inches(0.86), fill=TEAL, line=TEAL)
    add_textbox(slide, Inches(1.0), Inches(6.38), Inches(11.4), Inches(0.28), "The deck closes with a simple message: the product already reads as enterprise-grade, and the remaining work is focused on hardening, integration, and scale.", font_size=14, bold=True, color=NAVY)
    add_notes(slide, "End on the architecture and next phase. Make it clear what is already real and what comes next. Keep the tone confident: the product is already usable for demos and review.")
    add_footer(slide, 5, 5)

    out_path = OUT / "nexus-v2-presenter-deck.pptx"
    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    path = build()
    print(path)
